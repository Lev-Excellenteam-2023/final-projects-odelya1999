import argparse
import asyncio
import json
import openai
from pptx import Presentation


def extracting_the_text_from_a_single_slide(single_slide):
    """
        The function extracts the relevant content - the text
        :param single_slide: single slide of the powerpoint
        :return:The extracted text
        """
    final_text = " ".join([run.text for s in single_slide.shapes if s.has_text_frame
                           for single_paragraph in s.text_frame.paragraphs for run in single_paragraph.runs]).strip()
    return final_text


async def chat_with_gpt_3(slide_with_text):
    """
        The function adds the request to the text that we want to ask from gpt
        :param slide_with_text: single slide that contains fore sure a text
        :return: the response from gpt
        """
    final_req = slide_with_text + "Could you please explain to me what this text is about?"
    ans_from_gpt = await asyncio.get_event_loop().run_in_executor(None, lambda: openai.ChatCompletion.create
    (model="gpt-3.5-turbo", messages=[{"role": "user", "content": final_req}], timeout=100000))

    return ans_from_gpt.choices[0].message.content + "\n"


def creating_the_question_for_gpt(path_pp):
    """
        The function generates
        :param path_pp: The path the presentation is in
        :return: a list of tuples.
        """
    my_pp = Presentation(path_pp)
    tuples_list = []

    for i, s in enumerate(my_pp.slides, start=1):
        single_slide = extracting_the_text_from_a_single_slide(s)
        if single_slide:  # if the slide contains text.
            my_tuple = asyncio.create_task(chat_with_gpt_3(single_slide))
            tuples_list.append((i, single_slide, my_tuple))

    return tuples_list


async def put_the_answer_in_a_json_file(path_pp):
    """
       The function creates a json file for the answer from gpt
       :param path_pp: The path the presentation is in
       :return: the json file
       """
    tuples_list = creating_the_question_for_gpt(path_pp)
    response_dict = {}
    for index, slide_text, my_tuple in tuples_list:
        try:
            response = await my_tuple
        except Exception as e:
            response = f"Error occurred while processing slide {index}. Error message: {str(e)}"
        response_dict[f"response {index}"] = {"text": slide_text, "response": response}

    return json.dumps(response_dict, indent="\n")


async def main():
    """
    The main program.
    :return: present the result of program to the user.
    """
    my_path = argparse.ArgumentParser(description="Extracting the text from the PowerPoint presentation:")
    my_path.add_argument("File_pptx", type=str, help="The path that File_pptx is in")

    my_file = my_path.parse_args()
    name_of_file = my_file.File_pptx.split("/")[-1].split(".")[0]

    with open(f"{name_of_file}.json", "w") as fw:
        fw.write(await put_the_answer_in_a_json_file(my_file.File_pptx))

    print("The answer to your request is in the json file above :-)")

if __name__ == "_main_":
    openai.api_key = ""  # API key
    asyncio.get_event_loop().run_until_complete(main())